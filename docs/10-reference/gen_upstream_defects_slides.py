#!/usr/bin/env python3.9
"""Generate a 2-slide deck summarising genuine upstream XRootD findings.

Slide 1: Observed bugs (defects reproduced live).
Slide 2: Spec ambiguities / upstream quirks.

Run: python3.9 docs/10-reference/gen_upstream_defects_slides.py
Output: docs/10-reference/upstream-xrootd-defects.pptx
Source of truth: docs/10-reference/upstream-xrootd-defects.md
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- palette (dark deck) ---------------------------------------------------
BG = RGBColor(0x0F, 0x17, 0x2A)      # deep navy
PANEL = RGBColor(0x1B, 0x26, 0x3B)   # card
ACCENT = RGBColor(0x4F, 0xD1, 0xC5)  # teal (bugs)
ACCENT2 = RGBColor(0xF5, 0xB3, 0x4A)  # amber (quirks)
TEXT = RGBColor(0xE8, 0xED, 0xF6)
MUTED = RGBColor(0x9A, 0xA7, 0xBD)
CODE = RGBColor(0xC7, 0xF0, 0xEB)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _bg(slide):
    r = slide.shapes.add_shape(1, 0, 0, EMU_W, EMU_H)
    _fill(r, BG)
    r.shadow.inherit = False
    return r


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, space_after=4):
    """runs: list of paragraphs; each paragraph is list of (text,size,color,bold)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (t, sz, col, bold) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.color.rgb = col
            r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def _card(slide, left, top, width, height, accent):
    c = slide.shapes.add_shape(1, left, top, width, height)  # rounded rect
    try:
        c.adjustments[0] = 0.06
    except Exception:
        pass
    _fill(c, PANEL)
    c.shadow.inherit = False
    bar = slide.shapes.add_shape(1, left, top, Inches(0.09), height)
    _fill(bar, accent)
    bar.shadow.inherit = False
    return c


def header(slide, kicker, title, accent):
    _bg(slide)
    strip = slide.shapes.add_shape(1, 0, 0, EMU_W, Inches(0.14))
    _fill(strip, accent)
    strip.shadow.inherit = False
    _text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.4),
          [[(kicker, 14, accent, True)]])
    _text(slide, Inches(0.6), Inches(0.66), Inches(12.1), Inches(0.7),
          [[(title, 30, TEXT, True)]])


def bug_card(slide, left, top, w, h, tag, title, body, evidence, accent):
    _card(slide, left, top, w, h, accent)
    pad = Inches(0.26)
    _text(slide, left + pad, top + Inches(0.16), w - Inches(0.5), Inches(0.35),
          [[(tag + "  ", 12, accent, True), (title, 14.5, TEXT, True)]])
    _text(slide, left + pad, top + Inches(0.66), w - Inches(0.5), h - Inches(1.2),
          [[(body, 11.5, MUTED, False)]])
    _text(slide, left + pad, top + h - Inches(0.52), w - Inches(0.5), Inches(0.4),
          [[("● LIVE  ", 9.5, accent, True), (evidence, 9.5, CODE, False)]])


def build():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    # ================= SLIDE 1 — OBSERVED BUGS =================
    s1 = prs.slides.add_slide(blank)
    header(s1, "UPSTREAM XROOTD — DEVELOPMENT FINDINGS  |  1 of 2",
           "Confirmed upstream bugs (reproduced live)", ACCENT)

    col_w = Inches(4.06)
    gap = Inches(0.18)
    x0 = Inches(0.6)
    row_y = Inches(1.7)
    row_h = Inches(2.35)

    cards_top = [
        ("A1", "XrdCl sync-call deadlock",
         "A synchronous XrdCl call blocks forever in "
         "XrdSysCondVar::Wait with no observable timeout. GIL released, "
         "so even a SIGALRM watchdog can't fire — an un-interruptible hang.",
         "GDB backtrace; froze a 4,000-test run"),
        ("A2", "XrdCl dirlist framing corruption",
         "Under concurrent kXR_dirlist, a 2nd large response on a reused "
         "pooled connection desyncs the client parser → “Invalid response.” "
         "Module returned SUCCESS on the same load.",
         "Repeated across marathons, load-only"),
        ("A3", "Reference daemons die under load",
         "After a sustained marathon the stock xrootd/cmsd data-node "
         "daemons were found dead (proc count → 0); nginx workers survived. "
         "Thread-per-conn model meets fd/mem pressure.",
         "Observed at marathon end"),
    ]
    for i, (tag, title, body, ev) in enumerate(cards_top):
        bug_card(s1, x0 + i * (col_w + gap), row_y, col_w, row_h,
                 tag, title, body, ev, ACCENT)

    row2_y = row_y + row_h + Inches(0.2)
    row2_h = Inches(2.35)
    wide = Inches(6.2)
    cards_bot = [
        (x0, "A4", "CMS heartbeat drop → false NotFound",
         "A cmsd management link drops transiently under load while the "
         "data plane still serves bytes. Stock treats the node as gone and "
         "returns [3011] file not found for a file that is on disk.",
         "Deterministic: kill cmsd, data server up"),
        (x0 + wide + gap, "A5", "XrdHttp accepts non-conformant X.509",
         "Stock XrdHttp accepts out-of-namespace, wrong-CA-policy and "
         "CRL-revoked certs the WLCG/IGTF profile says to reject. TLS layer "
         "does no GSI/signing_policy verification at all.",
         "Differential + XRootD v6.1.0 source read"),
    ]
    for (x, tag, title, body, ev) in cards_bot:
        bug_card(s1, x, row2_y, wide, row2_h, tag, title, body, ev, ACCENT)

    _text(s1, Inches(0.6), Inches(7.02), Inches(12.1), Inches(0.35),
          [[("5 defects observed live — 2× XrdCl client, 2× server / CMS "
             "control plane, 1× XrdHttp security.  Ref: "
             "docs/11-architecture/reliability-under-load.md", 9.5, MUTED, False)]])

    # ================= SLIDE 2 — SPEC AMBIGUITIES =================
    s2 = prs.slides.add_slide(blank)
    header(s2, "UPSTREAM XROOTD — DEVELOPMENT FINDINGS  |  2 of 2",
           "Spec ambiguities & upstream quirks", ACCENT2)

    def quirk(slide, left, top, w, h, tag, title, body, note):
        _card(slide, left, top, w, h, ACCENT2)
        pad = Inches(0.26)
        _text(slide, left + pad, top + Inches(0.16), w - Inches(0.5),
              Inches(0.35),
              [[(tag + "  ", 12, ACCENT2, True), (title, 14.5, TEXT, True)]])
        _text(slide, left + pad, top + Inches(0.62), w - Inches(0.5),
              h - Inches(1.1), [[(body, 11.5, MUTED, False)]])
        _text(slide, left + pad, top + h - Inches(0.46), w - Inches(0.5),
              Inches(0.35), [[("◆ ", 9.5, ACCENT2, True), (note, 9.5, CODE, False)]])

    wide = Inches(6.2)
    q_y = Inches(1.7)
    q_h = Inches(2.5)
    quirk(s2, x0, q_y, wide, q_h, "B1",
          "Undocumented root:// wire behaviours",
          "Behaviours the official spec never describes: kXR_mv uses an "
          "ASCII-space separator (not NUL); real clients put one trailing NUL "
          "inside the path dlen; v5 sends handshake + kXR_protocol as one "
          "44-byte segment; kXR_new|kXR_delete means overwrite.",
          "Reverse-engineered from C++ + live clients")
    quirk(s2, x0 + wide + gap, q_y, wide, q_h, "B2",
          "xrdcp misleading 'key values mismatch'",
          "xrdcp prints ossl_x509_check_private_key: key values mismatch when "
          "handed an end-entity cert where it wants a proxy — even though the "
          "cert/key moduli match. A pure red herring.",
          "All 3 cert/key MD5s matched (ddc46a…)")

    q2_y = q_y + q_h + Inches(0.2)
    q2_h = Inches(2.15)
    narrow = Inches(4.06)
    trio = [
        ("B3", "mkdir 'idempotency'",
         "Stock returns rc=0 for a dir it created earlier in-process, but "
         "kXR_ItExists (3018) for a pre-existing on-disk dir. An XrdOss "
         "namespace-cache artifact, not a wire contract.",
         "Judged non-bug; not copied"),
        ("B4", "xrdgsiproxy info",
         "The stock tool ignores X509_USER_PROXY. This project honours it — "
         "the more useful behaviour — and records the stock behaviour as a "
         "quirk, a deliberate divergence.",
         "Deliberate divergence"),
        ("B5", "POSC disconnect handling",
         "On disconnect with an un-closed partial, stock keeps it pending a "
         "reconnect window; this module removes it immediately. A defensible "
         "semantic difference, left as a documented xfail.",
         "Defensible difference, not a bug"),
    ]
    for i, (tag, title, body, note) in enumerate(trio):
        quirk(s2, x0 + i * (narrow + gap), q2_y, narrow, q2_h,
              tag, title, body, note)

    _text(s2, Inches(0.6), Inches(7.02), Inches(12.1), Inches(0.35),
          [[("Quirks: the reference is self-consistent; the spec text (or a "
             "tool) is the awkward part.  Ref: docs/10-reference/"
             "protocol-notes.md · quirks.md", 9.5, MUTED, False)]])

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "upstream-xrootd-defects.pptx")
    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    build()
